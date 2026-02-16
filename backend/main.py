from fastapi import FastAPI, UploadFile, File, HTTPException, Form
from fastapi.responses import FileResponse, JSONResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from enum import Enum
import shutil
import os
import uuid
from pathlib import Path
import asyncio
from concurrent.futures import ThreadPoolExecutor
import io
import subprocess

# PyPDF2 imports
from PyPDF2 import PdfReader, PdfWriter

app = FastAPI(title="SimplePDF API", version="1.0.0")

# CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # 生产环境改具体域名
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 临时文件目录
TEMP_DIR = Path("/tmp/simplepdf")
TEMP_DIR.mkdir(exist_ok=True)

# 文件清理任务
async def cleanup_file(file_path: Path, delay: int = 300):
    """5分钟后自动清理文件"""
    await asyncio.sleep(delay)
    try:
        if file_path.exists():
            file_path.unlink()
            print(f"Cleaned up: {file_path}")
    except Exception as e:
        print(f"Cleanup error: {e}")

@app.post("/api/convert/pdf-to-excel")
async def convert_pdf_to_excel(file: UploadFile = File(...)):
    """PDF转Excel - 提取表格数据"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(400, "Only PDF files are allowed")
    
    # 验证文件大小 (50MB)
    content = await file.read()
    if len(content) > 50 * 1024 * 1024:
        raise HTTPException(400, "File too large. Max 50MB allowed.")
    
    import pdfplumber
    import openpyxl
    from openpyxl.styles import Font, Alignment, Border, Side
    
    file_id = str(uuid.uuid4())
    input_path = TEMP_DIR / f"{file_id}.pdf"
    output_path = TEMP_DIR / f"{file_id}.xlsx"
    
    try:
        # 保存上传文件
        with open(input_path, "wb") as f:
            f.write(content)
        
        tables_found = []
        all_tables = []
        
        # 打开PDF提取表格
        with pdfplumber.open(str(input_path)) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                tables = page.extract_tables()
                if tables:
                    for table_idx, table in enumerate(tables):
                        if table and len(table) > 0:
                            tables_found.append({
                                "page": page_num,
                                "table_index": table_idx + 1,
                                "rows": len(table),
                                "cols": len(table[0]) if table else 0
                            })
                            all_tables.append({
                                "page": page_num,
                                "table_index": table_idx + 1,
                                "data": table
                            })
        
        if not all_tables:
            raise HTTPException(400, "No tables found in PDF")
        
        # 创建Excel工作簿
        wb = openpyxl.Workbook()
        
        # 如果有多个表格，每个表格一个sheet
        # 如果只有一个表格，直接用默认sheet
        for idx, table_info in enumerate(all_tables):
            if idx == 0:
                ws = wb.active
                ws.title = f"Table_{table_info['page']}_{table_info['table_index']}"
            else:
                ws = wb.create_sheet(title=f"Table_{table_info['page']}_{table_info['table_index']}")
            
            # 写入表格数据
            for row_idx, row_data in enumerate(table_info['data'], 1):
                for col_idx, cell_value in enumerate(row_data, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=cell_value)
                    cell.alignment = Alignment(wrap_text=True, vertical='top')
            
            # 自动调整列宽
            for column in ws.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if cell.value:
                            cell_length = len(str(cell.value))
                            if cell_length > max_length:
                                max_length = min(cell_length, 50)  # 最大50字符
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column_letter].width = adjusted_width
        
        # 保存Excel文件
        wb.save(str(output_path))
        wb.close()
        
        print(f"Extracted {len(all_tables)} tables from PDF")
        print(f"Tables info: {tables_found}")
        
        # 启动清理任务
        asyncio.create_task(cleanup_file(input_path))
        asyncio.create_task(cleanup_file(output_path))
        
        return FileResponse(
            path=output_path,
            filename=file.filename.replace('.pdf', '.xlsx'),
            media_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            headers={
                "X-Tables-Count": str(len(all_tables)),
                "X-Tables-Info": str(tables_found)
            }
        )
        
    except HTTPException:
        raise
    except Exception as e:
        # 清理临时文件
        if input_path.exists():
            input_path.unlink()
        if output_path.exists():
            output_path.unlink()
        raise HTTPException(500, f"Conversion failed: {str(e)}")


@app.post("/api/convert/pdf-to-word")
async def convert_pdf_to_word(file: UploadFile = File(...)):
    """PDF转Word"""
    # 验证文件类型
    if not file.filename.endswith('.pdf'):
        raise HTTPException(400, "Only PDF files are allowed")
    
    # 验证文件大小 (50MB)
    content = await file.read()
    if len(content) > 50 * 1024 * 1024:
        raise HTTPException(400, "File too large. Max 50MB allowed.")
    
    # 生成唯一ID
    file_id = str(uuid.uuid4())
    input_path = TEMP_DIR / f"{file_id}.pdf"
    output_path = TEMP_DIR / f"{file_id}.docx"
    
    try:
        # 保存上传文件
        with open(input_path, "wb") as f:
            f.write(content)
        
        # 转换PDF到Word
        from pdf2docx import Converter
        cv = Converter(str(input_path))
        cv.convert(str(output_path), start=0, end=None)
        cv.close()
        
        # 启动清理任务
        asyncio.create_task(cleanup_file(input_path))
        asyncio.create_task(cleanup_file(output_path))
        
        return FileResponse(
            path=output_path,
            filename=file.filename.replace('.pdf', '.docx'),
            media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        )
        
    except Exception as e:
        # 清理临时文件
        if input_path.exists():
            input_path.unlink()
        if output_path.exists():
            output_path.unlink()
        raise HTTPException(500, f"Conversion failed: {str(e)}")

@app.post("/api/merge")
async def merge_pdfs(files: list[UploadFile] = File(...)):
    """合并多个PDF"""
    if len(files) < 2:
        raise HTTPException(400, "At least 2 files required")
    
    from PyPDF2 import PdfMerger
    
    file_id = str(uuid.uuid4())
    output_path = TEMP_DIR / f"merged_{file_id}.pdf"
    input_paths = []
    
    try:
        merger = PdfMerger()
        
        for file in files:
            if not file.filename.endswith('.pdf'):
                continue
            
            content = await file.read()
            input_path = TEMP_DIR / f"{uuid.uuid4()}.pdf"
            
            with open(input_path, "wb") as f:
                f.write(content)
            
            input_paths.append(input_path)
            merger.append(str(input_path))
        
        merger.write(str(output_path))
        merger.close()
        
        # 清理输入文件
        for path in input_paths:
            asyncio.create_task(cleanup_file(path, 0))
        asyncio.create_task(cleanup_file(output_path))
        
        return FileResponse(
            path=output_path,
            filename="merged.pdf",
            media_type='application/pdf'
        )
        
    except Exception as e:
        for path in input_paths:
            if path.exists():
                path.unlink()
        if output_path.exists():
            output_path.unlink()
        raise HTTPException(500, f"Merge failed: {str(e)}")

@app.post("/api/split")
async def split_pdf(
    file: UploadFile = File(...),
    pages: str = Form(...),  # 格式: "1,3,5-10" 或 "1-2,3-4,5-5"
):
    """拆分PDF - 返回ZIP压缩包"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(400, "Only PDF files allowed")
    
    from PyPDF2 import PdfReader, PdfWriter
    import zipfile
    
    file_id = str(uuid.uuid4())
    input_path = TEMP_DIR / f"{file_id}.pdf"
    output_dir = TEMP_DIR / f"split_{file_id}"
    output_dir.mkdir(exist_ok=True)
    
    try:
        content = await file.read()
        with open(input_path, "wb") as f:
            f.write(content)
        
        reader = PdfReader(str(input_path))
        
        # 解析页码范围
        print(f"Split request - Pages param: '{pages}', Total PDF pages: {len(reader.pages)}")
        
        if not pages or pages.strip() == '':
            raise HTTPException(400, f"No pages specified. Received: '{pages}'")
        
        # 解析每个分组
        groups = []  # 每个元素是一个元组 (group_name, [page_indices])
        group_index = 1
        
        for part in pages.split(','):
            part = part.strip()
            if not part:
                continue
            
            if '-' in part:
                # 范围格式: "1-2" -> 提取页面 0,1
                try:
                    start, end = map(int, part.split('-'))
                    if start < 1 or end < 1:
                        raise HTTPException(400, f"Page numbers must be >= 1: {part}")
                    if start > end:
                        raise HTTPException(400, f"Invalid range (start > end): {part}")
                    page_indices = list(range(start-1, end))
                    groups.append((f"pages_{start}-{end}", page_indices))
                except ValueError as e:
                    raise HTTPException(400, f"Invalid page range '{part}': {str(e)}")
            else:
                # 单页格式: "5" -> 提取页面 4
                try:
                    page_num = int(part)
                    if page_num < 1:
                        raise HTTPException(400, f"Page number must be >= 1: {part}")
                    groups.append((f"page_{page_num}", [page_num - 1]))
                except ValueError as e:
                    raise HTTPException(400, f"Invalid page number '{part}': {str(e)}")
        
        if not groups:
            raise HTTPException(400, "No valid pages to extract")
        
        print(f"Parsed {len(groups)} groups: {groups}")
        
        # 为每个组生成一个PDF
        generated_files = []
        for group_name, page_indices in groups:
            writer = PdfWriter()
            valid_pages = []
            
            for page_idx in page_indices:
                if 0 <= page_idx < len(reader.pages):
                    try:
                        writer.add_page(reader.pages[page_idx])
                        valid_pages.append(page_idx + 1)  # 转换为1-based页码
                    except Exception as e:
                        print(f"Error adding page {page_idx + 1}: {e}")
                else:
                    print(f"Warning: Page index {page_idx} out of range")
            
            if valid_pages:
                # 生成文件名
                if len(valid_pages) == 1:
                    output_filename = f"page_{valid_pages[0]}.pdf"
                else:
                    output_filename = f"pages_{valid_pages[0]}-{valid_pages[-1]}.pdf"
                
                output_path = output_dir / output_filename
                with open(output_path, "wb") as f:
                    writer.write(f)
                generated_files.append(output_path)
                print(f"Generated: {output_filename}")
        
        if not generated_files:
            raise HTTPException(400, "No pages could be extracted")
        
        # 打包成ZIP
        zip_path = TEMP_DIR / f"split_{file_id}.zip"
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for pdf_file in generated_files:
                zipf.write(pdf_file, pdf_file.name)
        
        print(f"Created ZIP: {zip_path} with {len(generated_files)} files")
        
        # 清理临时文件
        asyncio.create_task(cleanup_file(input_path))
        for pdf_file in generated_files:
            asyncio.create_task(cleanup_file(pdf_file))
        asyncio.create_task(cleanup_file(zip_path))
        
        # 清理目录
        import shutil
        if output_dir.exists():
            shutil.rmtree(output_dir)
        
        return FileResponse(
            path=zip_path,
            filename=f"split_{file.filename.replace('.pdf', '.zip')}",
            media_type='application/zip'
        )
        
    except HTTPException:
        raise
    except Exception as e:
        # 清理临时文件
        if input_path.exists():
            input_path.unlink()
        if output_dir.exists():
            import shutil
            shutil.rmtree(output_dir)
        raise HTTPException(500, f"Split failed: {str(e)}")

def compress_with_ghostscript(input_path: str, output_path: str, level: str) -> bool:
    """使用 Ghostscript 压缩 PDF，返回是否成功"""
    import shutil
    
    # 检查 Ghostscript 是否可用
    if not shutil.which("gs"):
        return False
    
    quality_settings = {
        "low": "/printer",
        "medium": "/ebook",
        "high": "/screen"
    }
    
    pdf_settings = quality_settings.get(level, "/ebook")
    
    cmd = [
        "gs",
        "-sDEVICE=pdfwrite",
        "-dCompatibilityLevel=1.4",
        f"-dPDFSETTINGS={pdf_settings}",
        "-dNOPAUSE",
        "-dQUIET",
        "-dBATCH",
        "-dColorImageDownsampleType=/Bicubic",
        "-dGrayImageDownsampleType=/Bicubic",
        "-dMonoImageDownsampleType=/Bicubic",
        f"-sOutputFile={output_path}",
        input_path
    ]
    
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        return result.returncode == 0 and os.path.exists(output_path)
    except Exception as e:
        print(f"Ghostscript error: {e}")
        return False

def compress_with_pypdf2(input_path: str, output_path: str) -> bool:
    """使用 PyPDF2 作为回退压缩方案"""
    try:
        reader = PdfReader(input_path)
        writer = PdfWriter()
        
        for page in reader.pages:
            writer.add_page(page)
        
        # 移除部分元数据以减少大小
        writer.add_metadata({"/Producer": "SimplePDF"})
        
        with open(output_path, "wb") as f:
            writer.write(f)
        
        return True
    except Exception as e:
        print(f"PyPDF2 compression error: {e}")
        return False

@app.post("/api/compress")
async def compress_pdf(
    file: UploadFile = File(...),
    level: str = Form("medium")  # low, medium, high
):
    """压缩PDF文件 - 优先使用 Ghostscript，回退到 PyPDF2"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(400, "Only PDF files allowed")
    
    file_id = str(uuid.uuid4())
    input_path = TEMP_DIR / f"{file_id}.pdf"
    output_path = TEMP_DIR / f"compressed_{file_id}.pdf"
    
    try:
        content = await file.read()
        original_size = len(content)
        
        with open(input_path, "wb") as f:
            f.write(content)
        
        # 优先尝试 Ghostscript
        gs_success = compress_with_ghostscript(str(input_path), str(output_path), level)
        
        if not gs_success:
            print("Ghostscript not available or failed, falling back to PyPDF2")
            pypdf_success = compress_with_pypdf2(str(input_path), str(output_path))
            
            if not pypdf_success:
                raise HTTPException(500, "Compression failed: both Ghostscript and PyPDF2 failed")
        
        compressed_size = output_path.stat().st_size
        compression_ratio = (1 - compressed_size / original_size) * 100
        
        method = "Ghostscript" if gs_success else "PyPDF2"
        print(f"{method} compression: {original_size} -> {compressed_size} ({compression_ratio:.1f}% reduction)")
        
        # 清理临时文件
        asyncio.create_task(cleanup_file(input_path))
        asyncio.create_task(cleanup_file(output_path))
        
        return FileResponse(
            path=output_path,
            filename=f"compressed_{file.filename}",
            media_type='application/pdf',
            headers={
                "X-Original-Size": str(original_size),
                "X-Compressed-Size": str(compressed_size),
                "X-Compression-Ratio": f"{compression_ratio:.1f}"
            }
        )
        
    except Exception as e:
        if input_path.exists():
            input_path.unlink()
        if output_path.exists():
            output_path.unlink()
        raise HTTPException(500, f"Compression failed: {str(e)}")

@app.post("/api/pdf-info")
async def get_pdf_info(file: UploadFile = File(...)):
    """获取PDF信息（页数等）"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(400, "Only PDF files allowed")
    
    from PyPDF2 import PdfReader
    
    file_id = str(uuid.uuid4())
    input_path = TEMP_DIR / f"{file_id}.pdf"
    
    try:
        content = await file.read()
        
        # 验证文件大小
        if len(content) > 50 * 1024 * 1024:
            raise HTTPException(400, "File too large. Max 50MB allowed.")
        
        with open(input_path, "wb") as f:
            f.write(content)
        
        reader = PdfReader(str(input_path))
        num_pages = len(reader.pages)
        
        # 立即清理临时文件
        if input_path.exists():
            input_path.unlink()
        
        return {
            "filename": file.filename,
            "pages": num_pages,
            "size": len(content)
        }
        
    except Exception as e:
        if input_path.exists():
            input_path.unlink()
        raise HTTPException(500, f"Failed to read PDF: {str(e)}")

@app.post("/api/convert/pdf-to-image")
async def convert_pdf_to_image(
    file: UploadFile = File(...),
    format: str = Form("jpg"),  # jpg, png
    dpi: int = Form(150)  # 分辨率
):
    """将 PDF 转换为图片（每页一张）"""
    print(f"pdf-to-image called: format={format}, dpi={dpi}, file={file.filename}")
    
    if not file.filename.endswith('.pdf'):
        raise HTTPException(400, "Only PDF files allowed")
    
    if format not in ["jpg", "jpeg", "png"]:
        raise HTTPException(400, "Format must be jpg or png")
    
    if dpi < 72 or dpi > 300:
        raise HTTPException(400, "DPI must be between 72 and 300")
    
    try:
        from pdf2image import convert_from_path
        import zipfile
        from PIL import Image
        print("pdf2image imported successfully")
    except ImportError as e:
        print(f"Import error: {e}")
        raise HTTPException(500, f"Missing dependency: {e}")
    
    file_id = str(uuid.uuid4())
    input_path = TEMP_DIR / f"{file_id}.pdf"
    output_dir = TEMP_DIR / f"images_{file_id}"
    
    try:
        output_dir.mkdir(exist_ok=True)
        content = await file.read()
        print(f"Received file: {len(content)} bytes")
        
        with open(input_path, "wb") as f:
            f.write(content)
        
        # 转换 PDF 为图片
        print(f"Converting PDF: {input_path}")
        ext = "png" if format == "png" else "jpeg"
        
        images = convert_from_path(
            str(input_path),
            dpi=dpi,
            fmt=ext,
            output_folder=str(output_dir),
            paths_only=True
        )
        
        print(f"Converted {len(images)} pages")
        
        image_files = [Path(img_path) for img_path in images]
        
        # 重命名文件
        for i, img_path in enumerate(image_files):
            new_name = output_dir / f"page_{i + 1}.{format}"
            img_path.rename(new_name)
            image_files[i] = new_name
            print(f"  Renamed to: {new_name}")
        
        if not image_files:
            raise HTTPException(400, "No pages could be converted")
        
        # 打包成 ZIP
        zip_path = TEMP_DIR / f"images_{file_id}.zip"
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for img_file in image_files:
                zipf.write(img_file, img_file.name)
        
        print(f"Created ZIP: {zip_path} with {len(image_files)} images")
        
        # 清理临时文件
        asyncio.create_task(cleanup_file(input_path))
        for img_file in image_files:
            asyncio.create_task(cleanup_file(img_file))
        asyncio.create_task(cleanup_file(zip_path))
        
        # 清理目录
        import shutil
        if output_dir.exists():
            shutil.rmtree(output_dir)
        
        return FileResponse(
            path=zip_path,
            filename=f"{file.filename.replace('.pdf', '')}_images.zip",
            media_type='application/zip',
            headers={
                "X-Total-Pages": str(len(image_files)),
                "X-Image-Format": format,
                "X-DPI": str(dpi)
            }
        )
        
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        error_detail = traceback.format_exc()
        print(f"PDF to Image ERROR: {error_detail}")
        # 清理临时文件
        if input_path.exists():
            input_path.unlink()
        if output_dir.exists():
            import shutil
            shutil.rmtree(output_dir)
        raise HTTPException(500, f"Conversion failed: {str(e)}")

@app.post("/api/convert/image-to-pdf")
async def convert_image_to_pdf(
    files: list[UploadFile] = File(...),
    page_size: str = Form("A4")  # A4, Letter, original
):
    """将图片转换为 PDF（支持多张图片合并）"""
    if not files:
        raise HTTPException(400, "No files uploaded")
    
    if len(files) > 20:
        raise HTTPException(400, "Max 20 images allowed")
    
    valid_extensions = {'.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.gif', '.webp'}
    
    for file in files:
        ext = Path(file.filename).suffix.lower()
        if ext not in valid_extensions:
            raise HTTPException(400, f"Invalid file type: {file.filename}. Only images allowed.")
    
    from PIL import Image
    
    file_id = str(uuid.uuid4())
    output_path = TEMP_DIR / f"images_{file_id}.pdf"
    images = []
    
    try:
        # 读取并转换所有图片
        for i, file in enumerate(files):
            content = await file.read()
            img = Image.open(io.BytesIO(content))
            
            # 转换为 RGB（PDF 不支持 RGBA）
            if img.mode == 'RGBA':
                # 创建白色背景
                background = Image.new('RGB', img.size, (255, 255, 255))
                background.paste(img, mask=img.split()[3])  # 使用 alpha 通道作为 mask
                img = background
            elif img.mode != 'RGB':
                img = img.convert('RGB')
            
            images.append(img)
            print(f"Processed image {i+1}: {file.filename}, mode: {img.mode}, size: {img.size}")
        
        if not images:
            raise HTTPException(400, "No valid images to convert")
        
        # 保存为 PDF
        first_image = images[0]
        remaining_images = images[1:] if len(images) > 1 else []
        
        first_image.save(
            output_path,
            "PDF",
            resolution=100.0,
            save_all=True,
            append_images=remaining_images
        )
        
        print(f"Created PDF: {output_path} with {len(images)} pages")
        
        # 清理临时文件
        for img in images:
            img.close()
        asyncio.create_task(cleanup_file(output_path))
        
        return FileResponse(
            path=output_path,
            filename=f"converted_images.pdf",
            media_type='application/pdf',
            headers={
                "X-Total-Images": str(len(images)),
                "X-Page-Size": page_size
            }
        )
        
    except HTTPException:
        raise
    except Exception as e:
        # 清理资源
        for img in images:
            try:
                img.close()
            except:
                pass
        if output_path.exists():
            output_path.unlink()
        raise HTTPException(500, f"Conversion failed: {str(e)}")

@app.get("/api/health")
async def health_check():
    """健康检查"""
    return {"status": "ok", "service": "simplepdf-api", "version": "1.0.1"}

# ==================== PDF to PPT ====================

@app.post("/api/convert/pdf-to-ppt")
async def convert_pdf_to_ppt(file: UploadFile = File(...)):
    """将 PDF 转换为 PPT，每页 PDF 转为 PPT 的一页"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(400, "Only PDF files are allowed")
    
    # 验证文件大小 (50MB)
    content = await file.read()
    if len(content) > 50 * 1024 * 1024:
        raise HTTPException(400, "File too large. Max 50MB allowed.")
    
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pdf2image import convert_from_path
    from PIL import Image
    import io
    
    file_id = str(uuid.uuid4())
    input_path = TEMP_DIR / f"{file_id}.pdf"
    output_path = TEMP_DIR / f"{file_id}.pptx"
    
    try:
        # 保存上传文件
        with open(input_path, "wb") as f:
            f.write(content)
        
        # 获取 PDF 页数
        reader = PdfReader(str(input_path))
        num_pages = len(reader.pages)
        
        if num_pages == 0:
            raise HTTPException(400, "PDF has no pages")
        
        # 将 PDF 转换为图片
        images = convert_from_path(str(input_path), dpi=150)
        
        # 创建 PPT
        prs = Presentation()
        
        # 设置幻灯片尺寸为 16:9 (默认是 4:3)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        for img in images:
            # 添加空白幻灯片
            blank_slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 将图片保存到内存
            img_bytes = io.BytesIO()
            img.save(img_bytes, format='PNG')
            img_bytes.seek(0)
            
            # 计算图片缩放比例以适应幻灯片
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            img_width, img_height = img.size
            aspect_ratio = img_width / img_height
            slide_ratio = slide_width / slide_height
            
            if aspect_ratio > slide_ratio:
                # 图片更宽，以宽度为准
                new_width = slide_width
                new_height = slide_width / aspect_ratio
                left = 0
                top = (slide_height - new_height) / 2
            else:
                # 图片更高，以高度为准
                new_height = slide_height
                new_width = slide_height * aspect_ratio
                left = (slide_width - new_width) / 2
                top = 0
            
            # 添加图片到幻灯片
            slide.shapes.add_picture(img_bytes, left, top, width=new_width, height=new_height)
        
        # 保存 PPT
        prs.save(str(output_path))
        
        # 启动清理任务
        asyncio.create_task(cleanup_file(input_path))
        asyncio.create_task(cleanup_file(output_path))
        
        return FileResponse(
            path=output_path,
            filename=file.filename.replace('.pdf', '.pptx'),
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            headers={"X-Total-Pages": str(num_pages)}
        )
        
    except HTTPException:
        raise
    except Exception as e:
        # 清理临时文件
        if input_path.exists():
            input_path.unlink()
        if output_path.exists():
            output_path.unlink()
        raise HTTPException(500, f"Conversion failed: {str(e)}")


# ==================== PDF Encrypt ====================

class Permission(str, Enum):
    """PDF 权限选项"""
    ALL = "all"  # 允许所有操作
    NO_PRINT = "no_print"  # 禁止打印
    NO_COPY = "no_copy"  # 禁止复制
    NO_EDIT = "no_edit"  # 禁止编辑
    NO_PRINT_COPY = "no_print_copy"  # 禁止打印和复制

@app.post("/api/protect/encrypt")
async def encrypt_pdf(
    file: UploadFile = File(...),
    password: str = Form(..., min_length=1),
    permission: Permission = Form(Permission.ALL)
):
    """为 PDF 添加密码保护"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(400, "Only PDF files are allowed")
    
    # 验证文件大小 (50MB)
    content = await file.read()
    if len(content) > 50 * 1024 * 1024:
        raise HTTPException(400, "File too large. Max 50MB allowed.")
    
    file_id = str(uuid.uuid4())
    input_path = TEMP_DIR / f"{file_id}.pdf"
    output_path = TEMP_DIR / f"encrypted_{file_id}.pdf"
    
    try:
        # 保存上传文件
        with open(input_path, "wb") as f:
            f.write(content)
        
        # 使用 pikepdf 进行加密
        import pikepdf
        
        pdf = pikepdf.open(str(input_path))
        
        # 设置权限标志
        # pikepdf 10.x Permissions 参数:
        # accessibility: 允许屏幕阅读器访问
        # extract: 允许提取文本/图形
        # modify_annotation: 允许修改注释
        # modify_assembly: 允许页面操作
        # modify_form: 允许填写表单
        # modify_other: 允许其他修改
        # print_lowres: 允许低分辨率打印
        # print_highres: 允许高分辨率打印
        
        if permission == Permission.ALL:
            # 允许所有操作
            allow_print = True
            allow_modify = True
            allow_copy = True
            allow_annot = True
        elif permission == Permission.NO_PRINT:
            allow_print = False
            allow_modify = True
            allow_copy = True
            allow_annot = True
        elif permission == Permission.NO_COPY:
            allow_print = True
            allow_modify = True
            allow_copy = False
            allow_annot = True
        elif permission == Permission.NO_EDIT:
            allow_print = True
            allow_modify = False
            allow_copy = True
            allow_annot = False
        elif permission == Permission.NO_PRINT_COPY:
            allow_print = False
            allow_modify = True
            allow_copy = False
            allow_annot = True
        else:
            allow_print = True
            allow_modify = True
            allow_copy = True
            allow_annot = True
        
        # 使用 R=6 (AES-256) 加密
        pdf.save(
            str(output_path),
            encryption=pikepdf.Encryption(
                owner=password,
                user=password,
                R=6,
                allow=pikepdf.Permissions(
                    accessibility=True,
                    extract=allow_copy,
                    modify_annotation=allow_annot,
                    modify_assembly=False,
                    modify_form=allow_annot,
                    modify_other=allow_modify,
                    print_lowres=allow_print,
                    print_highres=allow_print
                )
            )
        )
        pdf.close()
        
        # 启动清理任务
        asyncio.create_task(cleanup_file(input_path))
        asyncio.create_task(cleanup_file(output_path))
        
        return FileResponse(
            path=output_path,
            filename=f"encrypted_{file.filename}",
            media_type='application/pdf',
            headers={
                "X-Permission": permission.value,
                "X-Encrypted": "true"
            }
        )
        
    except HTTPException:
        raise
    except Exception as e:
        # 清理临时文件
        if input_path.exists():
            input_path.unlink()
        if output_path.exists():
            output_path.unlink()
        raise HTTPException(500, f"Encryption failed: {str(e)}")


# ==================== PDF Decrypt ====================

@app.post("/api/protect/decrypt")
async def decrypt_pdf(
    file: UploadFile = File(...),
    password: str = Form(..., min_length=1)
):
    """移除 PDF 密码保护"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(400, "Only PDF files are allowed")
    
    # 验证文件大小 (50MB)
    content = await file.read()
    if len(content) > 50 * 1024 * 1024:
        raise HTTPException(400, "File too large. Max 50MB allowed.")
    
    file_id = str(uuid.uuid4())
    input_path = TEMP_DIR / f"{file_id}.pdf"
    output_path = TEMP_DIR / f"decrypted_{file_id}.pdf"
    
    try:
        # 保存上传文件
        with open(input_path, "wb") as f:
            f.write(content)
        
        # 使用 pikepdf 进行解密
        import pikepdf
        
        try:
            # 尝试用密码打开
            pdf = pikepdf.open(str(input_path), password=password)
        except pikepdf._core.PasswordError:
            raise HTTPException(400, "Incorrect password")
        
        # 保存未加密的 PDF
        pdf.save(str(output_path))
        pdf.close()
        
        # 启动清理任务
        asyncio.create_task(cleanup_file(input_path))
        asyncio.create_task(cleanup_file(output_path))
        
        return FileResponse(
            path=output_path,
            filename=f"decrypted_{file.filename}",
            media_type='application/pdf',
            headers={"X-Decrypted": "true"}
        )
        
    except HTTPException:
        raise
    except Exception as e:
        # 清理临时文件
        if input_path.exists():
            input_path.unlink()
        if output_path.exists():
            output_path.unlink()
        raise HTTPException(500, f"Decryption failed: {str(e)}")


# ==================== PDF Watermark ====================

class WatermarkType(str, Enum):
    """水印类型"""
    TEXT = "text"
    IMAGE = "image"

class WatermarkPosition(str, Enum):
    """水印位置"""
    CENTER = "center"
    TOP_LEFT = "top-left"
    TOP_RIGHT = "top-right"
    BOTTOM_LEFT = "bottom-left"
    BOTTOM_RIGHT = "bottom-right"
    TILE = "tile"

@app.post("/api/watermark/add")
async def add_watermark(
    file: UploadFile = File(...),
    type: WatermarkType = Form(...),
    text: str = Form(None),
    image: UploadFile = File(None),
    position: WatermarkPosition = Form(WatermarkPosition.CENTER),
    opacity: float = Form(0.5),
    fontSize: int = Form(40),
    color: str = Form("#000000")
):
    """为 PDF 添加文字或图片水印"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(400, "Only PDF files are allowed")
    
    # 验证水印类型参数
    if type == WatermarkType.TEXT and not text:
        raise HTTPException(400, "Text is required for text watermark")
    if type == WatermarkType.IMAGE and not image:
        raise HTTPException(400, "Image is required for image watermark")
    
    # 验证透明度范围
    if opacity < 0 or opacity > 1:
        raise HTTPException(400, "Opacity must be between 0 and 1")
    
    # 验证文件大小 (50MB)
    content = await file.read()
    if len(content) > 50 * 1024 * 1024:
        raise HTTPException(400, "File too large. Max 50MB allowed.")
    
    from PIL import Image, ImageDraw, ImageFont
    from pdf2image import convert_from_path
    import img2pdf
    import io
    
    file_id = str(uuid.uuid4())
    input_path = TEMP_DIR / f"{file_id}.pdf"
    output_dir = TEMP_DIR / f"watermark_{file_id}"
    output_dir.mkdir(exist_ok=True)
    output_path = TEMP_DIR / f"watermarked_{file_id}.pdf"
    
    image_path = None
    
    try:
        # 保存上传文件
        with open(input_path, "wb") as f:
            f.write(content)
        
        # 获取 PDF 页数
        reader = PdfReader(str(input_path))
        num_pages = len(reader.pages)
        
        if num_pages == 0:
            raise HTTPException(400, "PDF has no pages")
        
        # 转换 PDF 为图片
        pdf_images = convert_from_path(str(input_path), dpi=150)
        
        watermarked_images = []
        
        for i, pdf_img in enumerate(pdf_images):
            # 转换为 RGBA 模式以支持透明度
            if pdf_img.mode != 'RGBA':
                pdf_img = pdf_img.convert('RGBA')
            
            # 创建透明层
            overlay = Image.new('RGBA', pdf_img.size, (255, 255, 255, 0))
            
            if type == WatermarkType.TEXT:
                draw = ImageDraw.Draw(overlay)
                
                # 解析颜色
                try:
                    r = int(color[1:3], 16)
                    g = int(color[3:5], 16)
                    b = int(color[5:7], 16)
                except:
                    r, g, b = 0, 0, 0
                
                # 计算字体大小 (基于 DPI 比例)
                font_size = int(fontSize * 2)  # 调整字体大小
                
                try:
                    font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", font_size)
                except:
                    font = ImageFont.load_default()
                
                # 获取文本尺寸
                bbox = draw.textbbox((0, 0), text, font=font)
                text_width = bbox[2] - bbox[0]
                text_height = bbox[3] - bbox[1]
                
                # 计算位置
                img_width, img_height = pdf_img.size
                
                if position == WatermarkPosition.CENTER:
                    x = (img_width - text_width) // 2
                    y = (img_height - text_height) // 2
                elif position == WatermarkPosition.TOP_LEFT:
                    x = 50
                    y = 50
                elif position == WatermarkPosition.TOP_RIGHT:
                    x = img_width - text_width - 50
                    y = 50
                elif position == WatermarkPosition.BOTTOM_LEFT:
                    x = 50
                    y = img_height - text_height - 50
                elif position == WatermarkPosition.BOTTOM_RIGHT:
                    x = img_width - text_width - 50
                    y = img_height - text_height - 50
                elif position == WatermarkPosition.TILE:
                    # 平铺模式
                    x_spacing = text_width + 100
                    y_spacing = text_height + 100
                    alpha = int(255 * opacity)
                    for row in range(0, img_height + y_spacing, y_spacing):
                        for col in range(0, img_width + x_spacing, x_spacing):
                            draw.text((col, row), text, font=font, fill=(r, g, b, alpha))
                    x, y = None, None  # 已经绘制完成
                else:
                    x = (img_width - text_width) // 2
                    y = (img_height - text_height) // 2
                
                # 绘制文字 (非平铺模式)
                if x is not None:
                    alpha = int(255 * opacity)
                    draw.text((x, y), text, font=font, fill=(r, g, b, alpha))
            
            else:  # IMAGE type
                # 如果是第一页，读取并保存水印图片
                if i == 0:
                    image_content = await image.read()
                    image_path = TEMP_DIR / f"watermark_img_{file_id}.png"
                    with open(image_path, "wb") as f:
                        f.write(image_content)
                
                # 打开水印图片
                watermark_img = Image.open(image_path)
                
                # 计算水印大小 (最大为页面 30%)
                img_width, img_height = pdf_img.size
                max_width = int(img_width * 0.3)
                max_height = int(img_height * 0.3)
                
                wm_width, wm_height = watermark_img.size
                ratio = min(max_width / wm_width, max_height / wm_height, 1.0)
                new_width = int(wm_width * ratio)
                new_height = int(wm_height * ratio)
                
                watermark_img = watermark_img.resize((new_width, new_height), Image.LANCZOS)
                
                # 转换为 RGBA
                if watermark_img.mode != 'RGBA':
                    watermark_img = watermark_img.convert('RGBA')
                
                # 调整透明度
                alpha = watermark_img.split()[-1]
                alpha = alpha.point(lambda p: int(p * opacity))
                watermark_img.putalpha(alpha)
                
                # 计算位置
                if position == WatermarkPosition.CENTER:
                    x = (img_width - new_width) // 2
                    y = (img_height - new_height) // 2
                elif position == WatermarkPosition.TOP_LEFT:
                    x = 50
                    y = 50
                elif position == WatermarkPosition.TOP_RIGHT:
                    x = img_width - new_width - 50
                    y = 50
                elif position == WatermarkPosition.BOTTOM_LEFT:
                    x = 50
                    y = img_height - new_height - 50
                elif position == WatermarkPosition.BOTTOM_RIGHT:
                    x = img_width - new_width - 50
                    y = img_height - new_height - 50
                elif position == WatermarkPosition.TILE:
                    # 平铺模式
                    x_spacing = new_width + 50
                    y_spacing = new_height + 50
                    for row in range(0, img_height + y_spacing, y_spacing):
                        for col in range(0, img_width + x_spacing, x_spacing):
                            overlay.paste(watermark_img, (col, row), watermark_img)
                    x, y = None, None
                else:
                    x = (img_width - new_width) // 2
                    y = (img_height - new_height) // 2
                
                # 粘贴水印 (非平铺模式)
                if x is not None:
                    overlay.paste(watermark_img, (x, y), watermark_img)
                
                watermark_img.close()
            
            # 合并图层
            result = Image.alpha_composite(pdf_img, overlay)
            result = result.convert('RGB')
            
            # 保存
            output_img_path = output_dir / f"page_{i+1}.jpg"
            result.save(output_img_path, 'JPEG', quality=95)
            watermarked_images.append(output_img_path)
            
            pdf_img.close()
            overlay.close()
        
        # 合并为 PDF
        with open(output_path, "wb") as f:
            f.write(img2pdf.convert([str(img) for img in watermarked_images]))
        
        # 清理临时文件
        asyncio.create_task(cleanup_file(input_path))
        for img_path in watermarked_images:
            asyncio.create_task(cleanup_file(img_path))
        asyncio.create_task(cleanup_file(output_path))
        
        # 清理目录
        import shutil
        if output_dir.exists():
            shutil.rmtree(output_dir)
        
        return FileResponse(
            path=output_path,
            filename=f"watermarked_{file.filename}",
            media_type='application/pdf',
            headers={
                "X-Watermark-Type": type.value,
                "X-Watermark-Position": position.value,
                "X-Total-Pages": str(num_pages)
            }
        )
        
    except HTTPException:
        raise
    except Exception as e:
        # 清理临时文件
        if input_path.exists():
            input_path.unlink()
        if output_dir.exists():
            import shutil
            shutil.rmtree(output_dir)
        if output_path.exists():
            output_path.unlink()
        if image_path and image_path.exists():
            image_path.unlink()
        raise HTTPException(500, f"Watermarking failed: {str(e)}")


if __name__ == "__main__":
    import uvicorn
    import os
    port = int(os.getenv("PORT", "8000"))
    uvicorn.run(app, host="0.0.0.0", port=port)
